from flask import Flask, render_template, request, session, redirect, send_from_directory
import os
from werkzeug.utils import secure_filename
from transformers import T5ForConditionalGeneration, T5Tokenizer, AutoModelForSeq2SeqLM, AutoTokenizer, pipeline
import nltk
from nltk.tokenize import sent_tokenize
from random import shuffle
from pptx import Presentation
from pdfminer.high_level import extract_text
from pathlib import Path
import shutil
from docx import Document
import logging
os.environ["HF_HUB_DISABLE_SYMLINKS_WARNING"] = "1"
logging.getLogger("transformers").setLevel(logging.ERROR)

# Download necessary NLTK data
nltk.download('punkt')
nltk.download('punkt_tab')

# Model identifiers
QG_MODEL_ID = "mrm8488/t5-base-finetuned-question-generation-ap"
QG_TOKENIZER_ID = "t5-base"
SUMMARIZER_ID = "sshleifer/distilbart-cnn-12-6"

# Local save paths
QG_MODEL_PATH = "models/qg-model"
QG_TOKENIZER_PATH = "models/qg-tokenizer"
SUM_MODEL_PATH = "models/summarizer"

# Hugging Face cache location (default)
HF_CACHE_DIR = os.path.expanduser("~/.cache/huggingface/hub")


def ensure_model(model_class, model_id, local_path):
    """Load a model from local path or download and save it there."""
    if os.path.exists(local_path):
        print(f"Found local model at {local_path}")
        return model_class.from_pretrained(local_path)
    else:
        print(f"Downloading model from Hugging Face: {model_id}")
        model = model_class.from_pretrained(model_id)
        model.save_pretrained(local_path)
        print(f"Saved to {local_path}")
        return model


def ensure_tokenizer(tokenizer_class, tokenizer_id, local_path):
    """Ensure a tokenizer exists locally or download/save it."""
    # Look for key tokenizer files
    required_files = ["tokenizer_config.json", "tokenizer.json", "vocab.json", "merges.txt"]
    files_exist = all(os.path.exists(os.path.join(local_path, f)) for f in required_files)

    if files_exist:
        print(f"Found local tokenizer at {local_path}")
        return tokenizer_class.from_pretrained(local_path)
    else:
        print(f"Downloading tokenizer from Hugging Face: {tokenizer_id}")
        tokenizer = tokenizer_class.from_pretrained(tokenizer_id)
        tokenizer.save_pretrained(local_path)
        print(f"Saved tokenizer to {local_path}")
        return tokenizer


def clear_huggingface_cache():
    """Optionally remove .cache files to save space after saving locally."""
    if os.path.exists(HF_CACHE_DIR):
        print(f"Deleting Hugging Face cache at {HF_CACHE_DIR}...")
        shutil.rmtree(HF_CACHE_DIR)
        print("Cache cleared.")
    else:
        print("No Hugging Face cache found.")


# ==== Load all models/tokenizers ====

# Question generation
qg_model = ensure_model(T5ForConditionalGeneration, QG_MODEL_ID, QG_MODEL_PATH)
qg_tokenizer = ensure_tokenizer(T5Tokenizer, QG_TOKENIZER_ID, QG_TOKENIZER_PATH)
generator = pipeline("text2text-generation", model=qg_model, tokenizer=qg_tokenizer)

# Summarization
sum_model = ensure_model(AutoModelForSeq2SeqLM, SUMMARIZER_ID, SUM_MODEL_PATH)
sum_tokenizer = ensure_tokenizer(AutoTokenizer, SUMMARIZER_ID, SUM_MODEL_PATH)
explainer = pipeline("summarization", model=sum_model, tokenizer=sum_tokenizer)

# === Optional: Delete cache to save space ===
clear_huggingface_cache()

# Flask app configuration
app = Flask(__name__)
app.secret_key = 'Notes2QA'
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER


def extract_text_from_pptx(path):
    """Extract text from a PPTX file."""
    prs = Presentation(path)
    full_text = []
    for slide in prs.slides:
        slide_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if slide_text:
            full_text.append(" ".join(slide_text))
    return "\n".join(full_text)


def extract_pdf_text(file_path):
    """Extract text from a PDF file."""
    return extract_text(file_path)


def is_valid_question(q):
    """Determine if a generated string is a valid question."""
    question_words = ('what', 'why', 'how', 'when', 'which', 'who', 'is', 'are', 'does', 'can')
    return q.endswith('?') and len(q.split()) > 5 and "?" in q and q.lower().startswith(question_words)


@app.route('/', methods=['GET', 'POST'])
def index():
    questions = []
    filename = None

    if request.method == 'POST':
        file = request.files.get('file')
        if file and file.filename:
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)
            session['file_path'] = file_path
        else:
            file_path = session.get('file_path')

        filename = os.path.basename(file_path)
        path = Path(f"explanations/{filename}")

        # Load from cache if exists
        if path.exists():
            with open(path, "r", encoding="utf-8") as f:
                chunks = f.read().strip().split("\n\n")
        else:
            # Extract data based on file type
            data = ""
            if file_path.lower().endswith(".pptx"):
                data = extract_text_from_pptx(file_path)
            elif file_path.lower().endswith(".docx"):
                doc = Document(file_path)
                data = "\n".join([para.text for para in doc.paragraphs])
            elif file_path.lower().endswith(".pdf"):
                data = extract_pdf_text(file_path)
            else:
                with open(file_path, "r", encoding="utf-8") as f:
                    data = f.read()

            sentences = sent_tokenize(data)
            num_sentences = len(sentences)
            print(f"{num_sentences} sentences")

            # Dynamically determine chunk size based on sentence count
            if num_sentences < 10:
                size = 1
            elif num_sentences < 50:
                size = 2
            elif num_sentences < 100:
                size = 4
            else:
                size = 6

            chunks = []
            for i in range(0, len(sentences), size):
                chunk = " ".join(sentences[i:i + size])
                chunk = explainer(chunk, max_length=100, min_length=30, do_sample=False)[0]['summary_text']
                chunks.append(chunk.strip())
                print(f"Explanation {i / size}:\n{chunk}")

            # Save chunks to cache
            with open(path, "w", encoding="utf-8") as f:
                for chunk in chunks:
                    f.write(chunk + "\n\n")

        shuffle(chunks)  # Randomize order

        # Generate up to 10 valid questions
        print("Generating questions...")
        questions = []
        used_chunks = set()
        for chunk in chunks:
            if len(questions) >= 10:
                break
            if chunk in used_chunks:
                continue

            input_text = f"Generate a question: {chunk}"
            result = generator(input_text, do_sample=False)
            q_text = result[0]['generated_text'].replace("question:", "").strip()

            if is_valid_question(q_text):
                questions.append(q_text)
                used_chunks.add(chunk)

        # Write generated questions to a text file
        path = Path("questions/questions.txt")
        with open(path, "w", encoding="utf-8") as f:
            for q in questions:
                f.write(q + "\n")

    return render_template('index.html', questions=questions, filename=filename)


@app.route('/clear_cache', methods=['POST'])
def clear_cache():
    """Clear cached explanation files."""
    cache_path = 'explanations/'

    if os.path.exists(cache_path):
        for filename in os.listdir(cache_path):
            file_path = os.path.join(cache_path, filename)
            try:
                if os.path.isfile(file_path):
                    os.remove(file_path)
                elif os.path.isdir(file_path):
                    shutil.rmtree(file_path)
            except Exception as e:
                print(f"Failed to delete {file_path}: {e}")
    return redirect('/')


@app.route('/download')
def download_questions():
    """Allow users to download the generated questions file."""
    questions_dir = os.path.join(os.getcwd(), "questions")
    questions_file = os.path.join(questions_dir, "questions.txt")

    if os.path.exists(questions_file):
        return send_from_directory(directory=questions_dir, path="questions.txt", as_attachment=True)
    return "Questions file not found", 404


if __name__ == '__main__':
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    app.run(debug=True)
